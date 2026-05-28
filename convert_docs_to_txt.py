#!/usr/bin/env python3
"""
批量将指定文件夹及其子文件夹中的 .doc, .docx, .ppt, .pptx 转换为 .txt
输出到另一个文件夹，保留原始目录结构
"""

import os
import sys
from pathlib import Path

# ==================== 配置区域 ====================
INPUT_DIR = "./docs_decomp"      # 原始文件所在目录
OUTPUT_DIR = "./docs_txt"        # 输出 TXT 的根目录
# =================================================

# 尝试导入必要的库
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("警告: python-docx 未安装，无法处理 .docx 文件")
    print("安装命令: pip install python-docx")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("警告: python-pptx 未安装，无法处理 .pptx 文件")
    print("安装命令: pip install python-pptx")

# 可选：处理旧版 .doc 和 .ppt 需要 textract（需额外安装）
try:
    import textract
    TEXTRACT_AVAILABLE = True
except ImportError:
    TEXTRACT_AVAILABLE = False
    print("提示: textract 未安装，旧版 .doc/.ppt 将无法处理")
    print("安装命令: pip install textract (可能需要额外系统依赖)")


def extract_docx(file_path: Path) -> str:
    """提取 .docx 文件中的文本"""
    doc = Document(file_path)
    text_parts = []
    
    # 提取段落
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    
    # 提取表格
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                text_parts.append(" | ".join(row_text))
    
    return "\n".join(text_parts)


def extract_pptx(file_path: Path) -> str:
    """提取 .pptx 文件中的文本"""
    prs = Presentation(file_path)
    text_parts = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        if slide_text:
            text_parts.append(f"[Slide {slide_idx}]")
            text_parts.extend(slide_text)
            text_parts.append("")  # 幻灯片间空行
    
    return "\n".join(text_parts)


def extract_with_textract(file_path: Path) -> str:
    """使用 textract 提取旧版 .doc/.ppt"""
    if not TEXTRACT_AVAILABLE:
        return ""
    try:
        text = textract.process(str(file_path)).decode("utf-8")
        return text
    except Exception as e:
        print(f"      textract 提取失败: {e}")
        return ""


def convert_file(input_path: Path, output_path: Path) -> bool:
    """
    转换单个文件，返回是否成功
    输出路径已在外部创建好目录
    """
    suffix = input_path.suffix.lower()
    
    # 如果输出文件已存在，跳过
    if output_path.exists():
        print(f"  跳过（已存在）: {output_path.name}")
        return True
    
    print(f"  转换: {input_path.name}")
    
    text = ""
    success = False
    
    if suffix == ".docx" and DOCX_AVAILABLE:
        try:
            text = extract_docx(input_path)
            success = True
        except Exception as e:
            print(f"    错误: {e}")
    
    elif suffix == ".pptx" and PPTX_AVAILABLE:
        try:
            text = extract_pptx(input_path)
            success = True
        except Exception as e:
            print(f"    错误: {e}")
    
    elif suffix in [".doc", ".ppt"]:
        text = extract_with_textract(input_path)
        if text:
            success = True
        else:
            print(f"    无法提取（需安装 textract 或手动转换）")
    
    else:
        print(f"    跳过不支持的类型: {suffix}")
        return False
    
    if success and text.strip():
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"    成功 → {output_path.name} ({len(text)} 字符)")
        return True
    elif success and not text.strip():
        print(f"    警告: 文件无文本内容")
        # 仍然创建空文件，避免重复处理
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("")
        return True
    
    return False


def main():
    # 检查依赖
    if not DOCX_AVAILABLE and not PPTX_AVAILABLE and not TEXTRACT_AVAILABLE:
        print("\n错误：没有可用的文档处理库，请先安装依赖")
        print("  pip install python-docx python-pptx")
        sys.exit(1)
    
    input_root = Path(INPUT_DIR)
    output_root = Path(OUTPUT_DIR)
    
    if not input_root.exists():
        print(f"错误：输入目录不存在: {INPUT_DIR}")
        sys.exit(1)
    
    # 收集所有需要转换的文件
    extensions = {".docx", ".pptx", ".doc", ".ppt"}
    files_to_convert = []
    
    for file_path in input_root.rglob("*"):
        if file_path.is_file() and file_path.suffix.lower() in extensions:
            files_to_convert.append(file_path)
    
    print(f"输入目录: {input_root}")
    print(f"输出目录: {output_root}")
    print(f"找到 {len(files_to_convert)} 个文件需要转换")
    print("-" * 50)
    
    # 转换文件
    success_count = 0
    for input_path in files_to_convert:
        # 计算相对路径，保留子文件夹结构
        rel_path = input_path.relative_to(input_root)
        # 替换扩展名为 .txt
        output_path = output_root / rel_path.with_suffix(".txt")
        # 确保输出目录存在
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        if convert_file(input_path, output_path):
            success_count += 1
    
    print("-" * 50)
    print(f"完成：成功 {success_count} / 总计 {len(files_to_convert)}")
    print(f"输出目录: {output_root.absolute()}")


if __name__ == "__main__":
    main()